import os
import sys

from pptx import Presentation

from pptx.dml.color import RGBColor
from pptx.util import Cm, Pt
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from pptx.enum.text import MSO_AUTO_SIZE

'''
    Sequence of slide layouts from python-pptx guide
    https://python-pptx.readthedocs.io/en/latest/user/slides.html

    Title (presentation title slide)
    Title and Content
    Section Header (sometimes called Segue)
    Two Content (side by side bullet textboxes)
    Comparison (same but additional title for each side by side content box)
    Title Only
    Blank
    Content with Caption
    Picture with Caption
'''

# Some util consants
STANZA = 'stanza'
TITLE = 'title'
IMAGE = 'image',
LOGO = 'logo'
FILENAME = 'filename'
MUSIC = 'music_letter'
PRESENTATION = 'presentation'


def usage():
    print(f'\nUtilização:\n  {sys.argv[0]} <<pasta_letras>> <<pasta_slides>> <<background_imagem>> <<logo_imagem>>\n')


def save_presentations_to_files(list_of_presentations, directory):
    previous_directory = os.getcwd()
    os.chdir(directory)
    file_extension = '.pptx'
    for presentation in list_of_presentations:
        final_filename = presentation[FILENAME][:-4] + file_extension
        final_presentation = presentation[PRESENTATION]
        final_presentation.save(final_filename)
    os.chdir(previous_directory)


def get_presentations_from_dictionary(list_of_dictionaries, background_image, logo_image):
    list_of_presentations = []
    for dictionary in list_of_dictionaries:
        title = dictionary[MUSIC][TITLE]
        stanzas = dictionary[MUSIC][STANZA]
        filename = dictionary[FILENAME]
        presentation = create_presentation_slide(title, stanzas, background_image, logo_image)
        list_of_presentations.append({
                                        FILENAME: filename,
                                        PRESENTATION: presentation
                                     })
    return list_of_presentations


def get_dictionaries_from_directory(directory):
    list_of_dictionaries = []
    previous_directory = os.getcwd()
    os.chdir(directory)
    list_of_filenames = os.listdir()
    for filename in list_of_filenames:
        music_letter_parts = read_file(filename)
        list_of_dictionaries.append({
                                        FILENAME: filename,
                                        MUSIC: music_letter_parts,
                                    })
    os.chdir(previous_directory)
    return list_of_dictionaries


def create_presentation_slide(title, stanzas, background_image, logo_image):
    presentation = Presentation()
    for stanza in stanzas:
        add_slide_with_title_and_content(presentation, background_image, logo_image, title, stanza)
    return presentation


def add_background_image(slide, image_path):
    top = Cm(0)
    left = Cm(0)
    height = Cm(19.05)
    width = Cm(25.40)
    slide.shapes.add_picture(image_file=image_path, left=left, top=top, height=height, width=width)


def add_title(slide, text):
    left = Cm(1.30)
    top = Cm(0.80)
    width = Cm(22.90)
    height = Cm(3.20)
    title = slide.shapes.add_textbox(left, top, width, height)
    title.text = text.upper()
    paragraph = title.text_frame.paragraphs[0]
    paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
    paragraph.font.size = Pt(40)
    paragraph.font.bold = True
    paragraph.font.italic = True
    paragraph.font.name = 'Comic Sans MS'
    paragraph.font.color.rgb = RGBColor(0, 176, 240)
    title.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    title.text_frame.word_wrap = True


def add_stanza(slide, text):
    left = Cm(1.20)
    top = Cm(5.80)
    width = Cm(23.00)
    height = Cm(12.00)
    stanza = slide.shapes.add_textbox(left, top, width, height)
    stanza.text = text.upper()
    stanza.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    stanza.text_frame.word_wrap = True
    for paragraph in stanza.text_frame.paragraphs:
        paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
        paragraph.font.size = Pt(36)
        paragraph.font.name = 'Comic Sans MS'
        paragraph.font.color.rgb = RGBColor(0, 176, 240)



def add_logo_bottom_right_corner(slide, image_path):
    top = Cm(15.60)
    left = Cm(20.10)
    height = Cm(3.10)
    width = Cm(5.00)
    slide.shapes.add_picture(image_file=image_path, left=left, top=top, height=height, width=width)


def add_slide_with_title_and_content(presentation, background_image_path, logo_image_path, title_text, stanza_text):
    blank_layout = presentation.slide_layouts[6]
    slide = presentation.slides.add_slide(blank_layout)

    add_background_image(slide, background_image_path)
    add_logo_bottom_right_corner(slide, logo_image_path)
    add_title(slide, title_text)
    add_stanza(slide, stanza_text)


def read_file(file_path):
    """
    :param file_path: Path of source music letter file
    :return: A dictionary containing the music title and a list of stanzas
    """
    dictionary = {}
    file = open(file_path, 'r')

    dictionary[TITLE] = file.readline()
    dictionary[STANZA] = []

    line = file.readline()
    stanza = ''
    while len(line) != 0:
        if line == '\n':
            if stanza != '':
                dictionary[STANZA].append(stanza)
            stanza = ''
        else:
            stanza = stanza + line
        line = file.readline()

    if stanza != '':
        dictionary[STANZA].append(stanza)

    file.close()
    return dictionary


if __name__ == '__main__':
    if len(sys.argv) != 5:
        usage()
        exit(0)

    directory = sys.argv[1]
    target_directory = sys.argv[2]
    background_image = sys.argv[3]
    logo_image = sys.argv[4]

    directory_dictionaries = get_dictionaries_from_directory(directory)
    list_of_presentations = get_presentations_from_dictionary(directory_dictionaries, background_image, logo_image)

    save_presentations_to_files(list_of_presentations, target_directory)
