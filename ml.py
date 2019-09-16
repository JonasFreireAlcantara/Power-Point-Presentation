from pptx.dml.color import RGBColor
from pptx.util import Cm, Pt
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from pptx.enum.text import MSO_AUTO_SIZE

'''
    Sequence of slide layouts from python-pptx guide
    https://python-pptx.readthedocs.io/en/latest/user/slides.html

    Title (presentation title slide)
    Title and Content
    Section Header (sometimes called Segue)
    Two Content (side by side bullet textboxes)
    Comparison (same but additional title for each side by side content box)
    Title Only
    Blank
    Content with Caption
    Picture with Caption
'''

# Some util consants
STANZA = 'STANZA'
TITLE = 'TITLE'


def add_background_image(slide, image_path):
    top = Cm(0)
    left = Cm(0)
    height = Cm(19.05)
    width = Cm(25.40)
    slide.shapes.add_picture(image_file=image_path, left=left, top=top, height=height, width=width)


def add_title(slide, text):
    left = Cm(1.30)
    top = Cm(0.80)
    width = Cm(22.90)
    height = Cm(3.20)
    title = slide.shapes.add_textbox(left, top, width, height)
    title.text = text.upper()
    paragraph = title.text_frame.paragraphs[0]
    paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
    paragraph.font.size = Pt(40)
    paragraph.font.bold = True
    paragraph.font.color.rgb = RGBColor(255, 255, 255)
    title.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE


def add_stanza(slide, text):
    left = Cm(1.20)
    top = Cm(5.80)
    width = Cm(23.00)
    height = Cm(12.00)
    stanza = slide.shapes.add_textbox(left, top, width, height)
    stanza.text = text.upper()
    for paragraph in stanza.text_frame.paragraphs:
        paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
        paragraph.font.size = Pt(30)
        paragraph.font.color.rgb = RGBColor(255, 255, 255)
    stanza.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE


def add_logo_bottom_right_corner(slide, image_path):
    top = Cm(14.65)
    left = Cm(21.60)
    height = Cm(4.40)
    width = Cm(3.80)
    slide.shapes.add_picture(image_file=image_path, left=left, top=top, height=height, width=width)


def add_slide_with_title_and_content(presentation, background_image_path, logo_image_path, title_text, stanza_text):
    blank_layout = presentation.slide_layouts[6]
    slide = presentation.slides.add_slide(blank_layout)

    add_background_image(slide, background_image_path)
    add_logo_bottom_right_corner(slide, logo_image_path)
    add_title(slide, title_text)
    add_stanza(slide, stanza_text)


def read_file(file_path):
    """
    :param file_path: Path of source music letter file
    :return: A dictionary containing the music title and a list of stanzas
    """
    dictionary = {}
    file = open(file_path, 'r')

    dictionary[TITLE] = file.readline()
    dictionary[STANZA] = []

    line = file.readline()
    stanza = ''
    while len(line) != 0:
        if line == '\n':
            if stanza != '':
                dictionary[STANZA].append(stanza)
            stanza = ''
        else:
            stanza = stanza + line
        line = file.readline()

    if stanza != '':
        dictionary[STANZA].append(stanza)
    return dictionary











    file.close()
